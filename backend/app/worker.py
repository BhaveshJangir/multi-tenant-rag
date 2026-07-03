import os
import pymupdf # PyMuPDF
import docx
import pptx
import openpyxl
from bs4 import BeautifulSoup
import mailparser
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from app.db.vector_store import qdrant_client, COLLECTION_NAME
from qdrant_client.http import models
import uuid
import asyncio

from app.core.config import settings

# Load the sentence transformer model
# Using the model configured in .env
embedding_model = SentenceTransformer(settings.EMBEDDING_MODEL_NAME)

def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    doc = pymupdf.open(file_path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path: str) -> str:
    prs = pptx.Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text.append(row_text)
    return "\n".join(text)

def extract_text_from_html(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
    return soup.get_text(separator="\n")

def extract_text_from_eml(file_path: str) -> str:
    mail = mailparser.parse_from_file(file_path)
    return mail.text_plain[0] if mail.text_plain else ""

def chunk_text(text: str) -> list[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    return splitter.split_text(text)

def process_document(document_id: str, file_path: str, tenant_id: str):
    """
    Background task to process an uploaded document.
    """
    try:
        # 1. Extract text
        ext = file_path.lower().split('.')[-1]
        if ext == 'pdf':
            text = extract_text_from_pdf(file_path)
        elif ext == 'docx':
            text = extract_text_from_docx(file_path)
        elif ext == 'pptx':
            text = extract_text_from_pptx(file_path)
        elif ext == 'xlsx':
            text = extract_text_from_xlsx(file_path)
        elif ext == 'html':
            text = extract_text_from_html(file_path)
        elif ext == 'eml':
            text = extract_text_from_eml(file_path)
        elif ext in ['csv', 'txt']:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            raise ValueError(f"Unsupported extension: {ext}")
                
        # 2. Chunk text
        chunks = chunk_text(text)
        
        # 3. Generate embeddings
        embeddings = embedding_model.encode(chunks)
        
        # 4. Index in Qdrant
        points = []
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
            points.append(
                models.PointStruct(
                    id=str(uuid.uuid4()),
                    vector=emb.tolist(),
                    payload={
                        "document_id": document_id,
                        "tenant_id": tenant_id,
                        "chunk_index": i,
                        "text": chunk
                    }
                )
            )
            
        qdrant_client.upsert(
            collection_name=COLLECTION_NAME,
            points=points
        )
        print(f"Successfully processed and indexed document {document_id}")
        
    except Exception as e:
        print(f"Error processing document {document_id}: {str(e)}")
        # Here we would update the document status to "error" in the DB
